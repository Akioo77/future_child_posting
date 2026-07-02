#!/usr/bin/env python3
"""Future Child Posting — 深度论文调研 PPT（详细版）"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_DARK   = RGBColor(0x1A, 0x2A, 0x3A)
C_PURPLE = RGBColor(0x6B, 0x4C, 0x9A)
C_PINK   = RGBColor(0xD4, 0x5B, 0x8B)
C_AMBER  = RGBColor(0xC8, 0x8B, 0x35)
C_BG     = RGBColor(0xF9, 0xF6, 0xFB)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xF0, 0xEB, 0xF5)
C_DARK2  = RGBColor(0x4A, 0x30, 0x6A)
C_GRAY   = RGBColor(0x7A, 0x6A, 0x8A)
C_GREEN  = RGBColor(0x2E, 0x8B, 0x6B)
C_TEAL   = RGBColor(0x00, 0x8B, 0x8B)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def bg(prs, color=C_BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color
    return slide

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, italic=False):
    txbox = slide.shapes.add_textbox(l, t, w, h)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txbox

def add_bullet(slide, items, l, t, w, h,
               font_size=12, color=C_DARK, indent=False):
    txbox = slide.shapes.add_textbox(l, t, w, h)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after  = Pt(2)
        prefix = "    • " if indent else "• "
        run = p.add_run()
        run.text = prefix + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txbox

def section_header(slide, title, subtitle="", color=C_PURPLE):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.8), color)
    add_text(slide, title, Inches(0.5), Inches(0.18), Inches(11), Inches(0.5),
             font_size=20, bold=True, color=C_WHITE)
    if subtitle:
        add_rect(slide, Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.45), C_PINK)
        add_text(slide, subtitle, Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.45),
                 font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    return slide

# ================================================================
def slide_cover(prs):
    slide = bg(prs, C_DARK2)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), C_PINK)
    add_text(slide, "Research Sharing", Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5),
             font_size=14, color=C_PINK)
    add_text(slide, "AI × Social Media",
             Inches(0.8), Inches(1.9), Inches(11.5), Inches(1.0),
             font_size=42, bold=True, color=C_WHITE)
    add_text(slide, "Future Child Posting",
             Inches(0.8), Inches(2.9), Inches(11.5), Inches(0.8),
             font_size=28, bold=True, color=C_PINK)
    add_text(slide, "家长在社交媒体分享孩子内容时的隐私风险与 AI 辅助保护研究",
             Inches(0.8), Inches(3.75), Inches(11.5), Inches(0.6),
             font_size=16, color=C_LIGHT)
    add_rect(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.05), C_PURPLE)
    add_text(slide, "三篇论文深度解析", Inches(0.8), Inches(4.7), Inches(11.5), Inches(0.4),
             font_size=13, color=C_GRAY)
    details = [
        "① Imago Obscura — 图片隐私 AI 助手（arXiv 2025，CHI 方向）",
        "② Parents & Children Privacy — 父母对儿童在线隐私的认知（2018）",
        "③ VLM Content Moderation — 多模态内容审核技术（arXiv 2023/2024）",
    ]
    add_bullet(slide, details, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.8),
               font_size=12, color=C_GRAY)
    add_rect(slide, Inches(10.0), Inches(5.6), Inches(3.0), Inches(1.5), C_PURPLE)
    add_text(slide, "📅 周末汇报", Inches(10.0), Inches(5.6), Inches(3.0), Inches(1.5),
             font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    return slide

# ================================================================
def slide_problem(prs):
    """问题定义 + Sharenting 现象"""
    slide = bg(prs, C_BG)
    section_header(slide, "问题背景  |  Sharenting 现象", "Background")

    # 左
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(5.9), Inches(6.1), C_WHITE, C_PURPLE, Pt(1.5))
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(5.9), Inches(0.5), C_PURPLE)
    add_text(slide, "什么是 Sharenting？", Inches(0.55), Inches(1.05), Inches(5.6), Inches(0.4),
             font_size=14, bold=True, color=C_WHITE)
    add_text(slide,
             "Sharenting = Sharing + Parenting\n"
             "指父母在社交媒体上主动分享孩子的生活信息。\n\n"
             "这一现象在 2015 年由《纽约时报》首次提出，"
             "现已发展为一个全球性的社会现象，并引发伦理讨论。",
             Inches(0.55), Inches(1.65), Inches(5.6), Inches(2.0),
             font_size=12, color=C_DARK)

    add_text(slide, "关键数据", Inches(0.55), Inches(3.7), Inches(5.6), Inches(0.35),
             font_size=13, bold=True, color=C_DARK2)
    stats = [
        ("140亿张/天", "全球每日社交媒体图片分享量（Broz, 2024）"),
        ("~2岁", "大多数孩子在此时已形成数字身份影子"),
        ("92%", "1-9 岁儿童家长曾在社交媒体分享孩子内容"),
        ("不可逆", "一旦发布，孩子终身无法撤回"),
    ]
    y = Inches(4.1)
    for val, desc in stats:
        add_rect(slide, Inches(0.55), y, Inches(1.8), Inches(0.55), C_LIGHT)
        add_text(slide, val, Inches(0.55), y + Inches(0.08), Inches(1.8), Inches(0.4),
                 font_size=13, bold=True, color=C_PURPLE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, Inches(2.45), y + Inches(0.1), Inches(3.7), Inches(0.45),
                 font_size=11, color=C_DARK)
        y += Inches(0.62)

    # 右
    add_rect(slide, Inches(6.6), Inches(1.0), Inches(6.4), Inches(6.1), C_WHITE, C_PURPLE, Pt(1.5))
    add_rect(slide, Inches(6.6), Inches(1.0), Inches(6.4), Inches(0.5), C_PURPLE)
    add_text(slide, "隐私风险分类（Imago Obscura Threat Model）", Inches(6.75), Inches(1.05), Inches(6.1), Inches(0.4),
             font_size=13, bold=True, color=C_WHITE)

    risks = [
        ("🔍 Observable Privacy（可观察的）",
         "图片中直接可见的敏感信息",
         "人脸、学校标志、车牌、门牌号、全名、生日、护照信息"),
        ("🧠 Inferential Privacy（可推断的）",
         "通过内容推断出的隐私信息",
         "位置习惯、社交关系、经济状况、健康信息、行为模式、宗教信仰"),
        ("📝 Contextual Privacy（情境的）",
         "来自发布情境的信息（非图片内容本身）",
         "图片说明文字、标签、@ 提及、地点打卡、发布时间规律"),
        ("⚖️ 伦理与 consent",
         "孩子没有同意权，终身影响",
         "尴尬未来、情感伤害、数字身份被商业化、无法撤回"),
    ]
    y = Inches(1.65)
    for title, subtitle, items in risks:
        add_rect(slide, Inches(6.7), y, Inches(6.2), Inches(1.45), C_LIGHT)
        add_text(slide, title, Inches(6.85), y + Inches(0.08), Inches(5.9), Inches(0.35),
                 font_size=12, bold=True, color=C_DARK2)
        add_text(slide, subtitle, Inches(6.85), y + Inches(0.42), Inches(5.9), Inches(0.28),
                 font_size=10, italic=True, color=C_PURPLE)
        add_text(slide, items, Inches(6.85), y + Inches(0.72), Inches(5.9), Inches(0.68),
                 font_size=10, color=C_GRAY)
        y += Inches(1.56)

    return slide

# ================================================================
def slide_papers_overview(prs):
    slide = bg(prs, C_BG)
    section_header(slide, "三篇论文总览  |  Papers Overview", "Overview")

    papers = [
        {
            "num": "①",
            "title": "Imago Obscura",
            "venue": "arXiv 2025 · CHI 方向 · cs.HC",
            "authors": "Kyzyl Monteiro, Yuchen Wu, Sauvik Das（美国）",
            "focus": "图片隐私 AI Copilot",
            "method": "Formative Study（7名专家）→ 系统开发 → User Study（15人）",
            "key": "5个设计要求 + 5类风险 + 8种脱敏技术",
            "color": C_PURPLE,
            "tag": "核心参考",
        },
        {
            "num": "②",
            "title": "Parents & Children Online Privacy",
            "venue": "Proc. Privacy & Trust 2018",
            "authors": "Jun Zhao",
            "focus": "父母对儿童在线隐私风险的认知与应对",
            "method": "访谈 + 问卷调查（6-11岁儿童家长）",
            "key": "父母认知偏差：低估风险 + 保护性策略为主",
            "color": C_AMBER,
            "tag": "用户理解",
        },
        {
            "num": "③",
            "title": "VLM Content Moderation",
            "venue": "arXiv 2023/2024",
            "authors": "Ahmed, Hu, Sukthankar",
            "focus": "VLM 零样本 + 多模态融合的儿童内容审核",
            "method": "技术评估 + 零样本实验 + 音频视觉融合",
            "key": "自然语言提示指导 VLM 识别细微不当内容",
            "color": C_PINK,
            "tag": "技术路线",
        },
    ]

    x = Inches(0.35)
    for p in papers:
        w = Inches(4.18)
        add_rect(slide, x, Inches(1.0), w, Inches(6.1), C_WHITE, p["color"], Pt(2))
        add_rect(slide, x, Inches(1.0), w, Inches(0.55), p["color"])
        add_text(slide, p["num"], x + Inches(0.12), Inches(1.07), Inches(0.45), Inches(0.42),
                 font_size=16, bold=True, color=C_WHITE)
        add_rect(slide, x + w - Inches(1.2), Inches(1.1), Inches(1.1), Inches(0.35), p["color"])
        add_text(slide, p["tag"], x + w - Inches(1.2), Inches(1.1), Inches(1.1), Inches(0.35),
                 font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        add_text(slide, p["title"], x + Inches(0.15), Inches(1.68), w - Inches(0.3), Inches(0.7),
                 font_size=13, bold=True, color=C_DARK2)
        add_text(slide, p["venue"], x + Inches(0.15), Inches(2.38), w - Inches(0.3), Inches(0.35),
                 font_size=10, color=C_GRAY)
        add_text(slide, p["authors"], x + Inches(0.15), Inches(2.73), w - Inches(0.3), Inches(0.35),
                 font_size=10, color=p["color"])

        add_rect(slide, x + Inches(0.15), Inches(3.13), w - Inches(0.3), Inches(0.03), p["color"])

        add_text(slide, "研究焦点", x + Inches(0.15), Inches(3.23), w - Inches(0.3), Inches(0.28),
                 font_size=10, bold=True, color=C_GRAY)
        add_text(slide, p["focus"], x + Inches(0.15), Inches(3.51), w - Inches(0.3), Inches(0.65),
                 font_size=11, color=C_DARK)

        add_text(slide, "研究方法", x + Inches(0.15), Inches(4.23), w - Inches(0.3), Inches(0.28),
                 font_size=10, bold=True, color=C_GRAY)
        add_text(slide, p["method"], x + Inches(0.15), Inches(4.51), w - Inches(0.3), Inches(0.85),
                 font_size=11, color=C_DARK)

        add_text(slide, "核心贡献", x + Inches(0.15), Inches(5.43), w - Inches(0.3), Inches(0.28),
                 font_size=10, bold=True, color=C_GRAY)
        add_text(slide, p["key"], x + Inches(0.15), Inches(5.71), w - Inches(0.3), Inches(1.0),
                 font_size=11, color=C_DARK2)

        x += Inches(4.43)

    return slide

# ================================================================
def slide_imago_threat(prs):
    slide = bg(prs, C_BG)
    """Imago Obscura — Threat Model + 5类风险"""
    section_header(slide, "论文 1  |  Imago Obscura — Threat Model & 五大风险类别", "① 核心参考")

    # 左：Threat Model
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(4.1), Inches(6.1), C_WHITE, C_PURPLE, Pt(1.5))
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(4.1), Inches(0.5), C_PURPLE)
    add_text(slide, "Threat Model（威胁模型）", Inches(0.55), Inches(1.05), Inches(3.9), Inches(0.4),
             font_size=13, bold=True, color=C_WHITE)
    add_text(slide, "对手目标：", Inches(0.55), Inches(1.62), Inches(3.9), Inches(0.3),
             font_size=11, bold=True, color=C_DARK2)
    add_text(slide, "从图片的边缘信息（非分享意图核心内容）中推断敏感个人信息",
             Inches(0.55), Inches(1.92), Inches(3.9), Inches(0.65),
             font_size=11, color=C_DARK)
    add_text(slide, "覆盖范围", Inches(0.55), Inches(2.65), Inches(3.9), Inches(0.3),
             font_size=11, bold=True, color=C_DARK2)
    scope = [
        "✅ Observable Privacy（可见信息，如人脸、车牌）",
        "✅ Inferential Privacy（可推断信息，如位置习惯）",
        "❌ Contextual Privacy（文字说明、metadata）不纳入",
        "❌ 算法对抗攻击（如对抗性扰动）不纳入",
    ]
    add_bullet(slide, scope, Inches(0.55), Inches(2.95), Inches(3.9), Inches(1.9),
               font_size=10, color=C_DARK)

    add_text(slide, "隐私风险的后果", Inches(0.55), Inches(4.9), Inches(3.9), Inches(0.3),
             font_size=11, bold=True, color=C_DARK2)
    harms = ["尴尬 / Embarrassment", "失业 / Job loss",
             "身份盗窃 / Identity theft", "跟踪 / Stalking", "骚扰 / Harassment"]
    add_bullet(slide, harms, Inches(0.55), Inches(5.2), Inches(3.9), Inches(1.8),
               font_size=10, color=C_RED)

    # 右：5类风险
    add_rect(slide, Inches(4.75), Inches(1.0), Inches(8.15), Inches(6.1), C_WHITE, C_PURPLE, Pt(1.5))
    add_rect(slide, Inches(4.75), Inches(1.0), Inches(8.15), Inches(0.5), C_PURPLE)
    add_text(slide, "5 类图片隐私风险（Imago Obscura 识别范围）", Inches(4.9), Inches(1.05), Inches(7.8), Inches(0.4),
             font_size=13, bold=True, color=C_WHITE)

    five_risks = [
        ("🔐 Self-Disclosure Risk", "自我披露风险",
         "意外暴露个人习惯、健康状况、生活事件",
         "例：生日蛋糕上的数字蜡烛 → 暴露孩子年龄\n书架照片 → 推断政治/知识倾向\n药瓶 → 暴露健康问题"),
        ("🪪 Identity Exposure Risk", "身份暴露风险",
         "通过可见的身份信息识别个人",
         "例：纹身（独特的纹身可定位个人）\n虹膜 / 面部特征（即使脸被遮挡）\nID 卡、邮件上的姓名地址"),
        ("📋 Confidential Info Leakage", "机密信息泄露",
         "背景中可见的商业/私人敏感信息",
         "例：白板上的项目计划 → 商业机密\n电脑屏幕内容 → 账号密码\n桌面文件 → 私人地址"),
        ("📍 Location Exposure Risk", "位置暴露风险",
         "可识别地理位置的建筑/地标/天气",
         "例：窗外可见的城市天际线 → 定位住址\n学校标志 → 孩子就读学校\n标志性建筑 → 常去地点"),
        ("👤 Bystander Risk", "旁观者风险",
         "背景中非自愿被拍摄的路人隐私",
         "例：街上路人被摄入\n学校活动中的其他孩子\n商场/公园的旁观者"),
    ]

    y = Inches(1.62)
    for icon_title, cn_title, desc_en, examples in five_risks:
        add_rect(slide, Inches(4.85), y, Inches(7.95), Inches(1.1), C_LIGHT)
        add_text(slide, f"{icon_title} — {cn_title}", Inches(4.98), y + Inches(0.06), Inches(7.7), Inches(0.32),
                 font_size=11, bold=True, color=C_PURPLE)
        add_text(slide, desc_en, Inches(4.98), y + Inches(0.38), Inches(7.7), Inches(0.28),
                 font_size=10, italic=True, color=C_GRAY)
        add_text(slide, examples, Inches(4.98), y + Inches(0.66), Inches(7.7), Inches(0.4),
                 font_size=9, color=C_DARK)
        y += Inches(1.18)

    return slide

# ================================================================
def slide_imago_dr(prs):
    slide = bg(prs, C_BG)
    """Imago Obscura — 5个设计要求详解"""
    section_header(slide, "论文 1  |  Imago Obscura — 五大设计要求（DR1-5）", "① 核心参考")

    drs = [
        {
            "id": "DR1",
            "title": "支持表达分享意图和隐私担忧",
            "detail": "参与者选择同一张图片时，往往标记不同的敏感区域——这取决于他们想发给谁、为什么发。系统应适配用户直接表达的隐私担忧和分享意图。",
            "example": "用户输入：'我想发孩子生日照片，但不想露脸' → 系统理解意图后自动聚焦面部区域",
            "color": C_PURPLE,
        },
        {
            "id": "DR2",
            "title": "提高内容层面的隐私风险意识",
            "detail": "许多用户不了解某些隐私风险类型。系统应在用户直接担忧的基础上，主动提示用户可能忽略的风险（超过用户即时意识范围的风险）。",
            "example": "用户只担心露脸 → 系统还提示：照片背景中的学校标志也属于位置暴露风险",
            "color": C_TEAL,
        },
        {
            "id": "DR3",
            "title": "促进明智决策",
            "detail": "目标是最大程度降低隐私风险，同时最小程度影响分享意图。系统应解释风险含义，推荐对分享意图影响最小的脱敏策略。",
            "example": "将「模糊」和「生成替换」两种方案的效果并排展示，让用户权衡",
            "color": C_AMBER,
        },
        {
            "id": "DR4",
            "title": "便于有效应用脱敏技术",
            "detail": "即使是专业图像编辑者也只会使用最熟悉的脱敏技术（而非最有效的）。系统应简化操作流程，降低精确选择敏感内容和有效应用脱敏技术的门槛。",
            "example": "一键选择对象 + 一键应用推荐脱敏方案，无需手动繁琐操作",
            "color": C_GREEN,
        },
        {
            "id": "DR5",
            "title": "确保自主性和精细控制",
            "detail": "用户在乎的不仅是脱敏，还有图片真实性和美观度。用户偏好因隐私担忧、分享意图和个人审美而异。系统应只是一个 Copilot，最终决定权在用户。",
            "example": "用户可以手动调整每个脱敏区域的范围和强度，AI 只推荐，用户拍板",
            "color": C_PINK,
        },
    ]

    y = Inches(1.0)
    for dr in drs:
        h = Inches(1.18)
        add_rect(slide, Inches(0.4), y, Inches(12.5), h, C_WHITE, dr["color"], Pt(1.5))
        add_rect(slide, Inches(0.4), y, Inches(0.65), h, dr["color"])
        add_text(slide, dr["id"], Inches(0.4), y, Inches(0.65), h,
                 font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, dr["title"], Inches(1.15), y + Inches(0.08), Inches(11.5), Inches(0.38),
                 font_size=12, bold=True, color=C_DARK2)
        add_text(slide, dr["detail"], Inches(1.15), y + Inches(0.46), Inches(7.5), Inches(0.65),
                 font_size=10, color=C_DARK)
        add_rect(slide, Inches(8.75), y + Inches(0.12), Inches(4.0), Inches(0.9), RGBColor(0xF5, 0xF3, 0xF8))
        add_text(slide, "💡 " + dr["example"], Inches(8.88), y + Inches(0.2), Inches(3.8), Inches(0.8),
                 font_size=10, color=C_DARK2, italic=True)
        y += Inches(1.26)

    return slide

# ================================================================
def slide_imago_tech(prs):
    slide = bg(prs, C_BG)
    """Imago Obscura — 技术实现细节"""
    section_header(slide, "论文 1  |  Imago Obscura — 技术实现 & 8种脱敏技术", "① 核心参考")

    # 左：系统架构
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(5.9), Inches(6.1), C_WHITE, C_PURPLE, Pt(1.5))
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(5.9), Inches(0.5), C_PURPLE)
    add_text(slide, "系统架构（AI 模型组合）", Inches(0.55), Inches(1.05), Inches(5.6), Inches(0.4),
             font_size=13, bold=True, color=C_WHITE)

    pipeline = [
        ("Step 1", "📷 图片上传", "用户上传待分析图片", C_LIGHT),
        ("Step 2", "🔍 Vision Model", "识别并标注图片中所有对象\n(Xiao et al., 2023 — Grounding DINO)", RGBColor(0xE8, 0xE0, 0xF5)),
        ("Step 3", "🧠 MLLM（多模态大模型）", "根据用户描述的担忧，识别相关隐私风险\n(GPT-4V 或等效模型)", C_LIGHT),
        ("Step 4", "✂️ Segmentation Model", "精确分割需要脱敏的区域\n(SAM — Segment Anything, Kirillov 2023)", RGBColor(0xE8, 0xE0, 0xF5)),
        ("Step 5", "🎨 Image Generation", "应用脱敏技术：Inpainting / 模糊 / 生成替换\n(Stable Diffusion / StyleGAN)", C_LIGHT),
        ("Step 6", "✅ 输出", "修改后图片 + 风险说明 → 用户决策", RGBColor(0xE8, 0xE0, 0xF5)),
    ]
    y = Inches(1.62)
    for step, title, desc, col in pipeline:
        add_rect(slide, Inches(0.55), y, Inches(5.65), Inches(0.88), col, C_PURPLE, Pt(1))
        add_text(slide, step, Inches(0.62), y + Inches(0.05), Inches(0.55), Inches(0.3),
                 font_size=9, bold=True, color=C_PURPLE)
        add_text(slide, title, Inches(1.2), y + Inches(0.06), Inches(4.9), Inches(0.3),
                 font_size=11, bold=True, color=C_DARK2)
        add_text(slide, desc, Inches(1.2), y + Inches(0.4), Inches(4.9), Inches(0.45),
                 font_size=10, color=C_GRAY)
        y += Inches(0.95)

    # 右：8种脱敏技术
    add_rect(slide, Inches(6.6), Inches(1.0), Inches(6.4), Inches(6.1), C_WHITE, C_PURPLE, Pt(1.5))
    add_rect(slide, Inches(6.6), Inches(1.0), Inches(6.4), Inches(0.5), C_PURPLE)
    add_text(slide, "8 种图片脱敏技术（来自 Li et al., 2017b 等研究）", Inches(6.75), Inches(1.05), Inches(6.1), Inches(0.4),
             font_size=12, bold=True, color=C_WHITE)

    techniques = [
        ("Pixelation（像素化）", "马赛克", "低", "低", "最常用，但可被部分还原"),
        ("Blurring（模糊）", "高斯模糊", "中", "低", "适合背景，使用最广"),
        ("Black Bar（黑条遮挡）", "黑色矩形", "高", "中", "直接但影响美观"),
        ("Cropping（裁剪）", "裁掉敏感区域", "高", "高", "最干净，但丢失图片内容"),
        ("Inpainting（智能修补）", "AI 填充背景", "高", "高", "自然，但生成有随机性"),
        ("Avatar Replacement（虚拟形象）", "替换为 AI 头像", "高", "高", "保护身份，保持趣味性"),
        ("Generative Replace（生成替换）", "AI 生成替代内容", "高", "高", "最自然，像「魔法编辑」"),
        ("Silhouette（剪影）", "转为黑色轮廓", "高", "中", "保留姿态信息，隐藏身份"),
    ]

    # 表头
    add_rect(slide, Inches(6.65), Inches(1.62), Inches(6.3), Inches(0.35), C_LIGHT)
    headers = [("技术名称", 1.8), ("中文名", 1.0), ("有效性", 0.65), ("美观度", 0.65), ("特点", 2.2)]
    x = Inches(6.7)
    for hdr, w in headers:
        add_text(slide, hdr, x, Inches(1.65), Inches(w), Inches(0.3),
                 font_size=9, bold=True, color=C_DARK2, align=PP_ALIGN.CENTER)
        x += Inches(w)

    y = Inches(2.0)
    for name, cn, eff, beauty, note in techniques:
        bg_col = C_LIGHT if techniques.index((name, cn, eff, beauty, note)) % 2 == 0 else C_WHITE
        add_rect(slide, Inches(6.65), y, Inches(6.3), Inches(0.5), bg_col)
        x = Inches(6.7)
        cells = [(name, 1.8, C_DARK), (cn, 1.0, C_GRAY), (eff, 0.65, C_TEAL if eff == "高" else C_AMBER),
                 (beauty, 0.65, C_TEAL if beauty == "高" else C_AMBER), (note, 2.2, C_GRAY)]
        for val, w, col in cells:
            add_text(slide, val, x, y + Inches(0.08), Inches(w), Inches(0.38),
                     font_size=9, color=col)
            x += Inches(w)
        y += Inches(0.52)

    add_text(slide, "注：有效性 = 对身份/信息识别的保护程度；美观度 = 保持图片自然感的程度",
             Inches(6.65), Inches(6.55), Inches(6.3), Inches(0.4),
             font_size=9, italic=True, color=C_GRAY)

    return slide

# ================================================================
def slide_imago_user_study(prs):
    slide = bg(prs, C_BG)
    """Imago Obscura — Formative Study 细节 + User Study 结果"""
    section_header(slide, "论文 1  |  Imago Obscura — Formative Study & User Study 详解", "① 核心参考")

    # 左：Formative Study
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(6.0), Inches(6.1), C_WHITE, C_PURPLE, Pt(1.5))
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(6.0), Inches(0.5), C_PURPLE)
    add_text(slide, "Formative Study（7名图像编辑专家）", Inches(0.55), Inches(1.05), Inches(5.7), Inches(0.4),
             font_size=13, bold=True, color=C_WHITE)

    add_text(slide, "研究设计", Inches(0.55), Inches(1.62), Inches(5.7), Inches(0.3),
             font_size=11, bold=True, color=C_DARK2)
    procedure = [
        "数据来源：DIPA 数据集（115+ 张图片，带敏感内容标注）",
        "任务：选择图片 → 设想分享意图和隐私担忧 → 在 Krita 中脱敏",
        "提供的 AI 工具：对象分割、边界框、生成替换、参考图生成、头像替换",
        "还提供了隐私知识材料：敏感内容列表、威胁类型、脱敏前后对比示例",
        "方法：图像脱敏任务 + 任务后半结构化访谈 + 屏幕录制",
    ]
    add_bullet(slide, procedure, Inches(0.55), Inches(1.92), Inches(5.7), Inches(2.1),
               font_size=10, color=C_DARK)

    add_text(slide, "关键发现（Pain Points）", Inches(0.55), Inches(4.1), Inches(5.7), Inches(0.3),
             font_size=11, bold=True, color=C_DARK2)
    pain_pts = [
        "P1：同一张图片，不同人标记不同敏感区域（取决于分享意图）",
        "P2：用户参照隐私列表，但觉得太长太复杂（'看一遍就不想再看'）",
        "P3：即使专家也只选自己熟悉的技术，而非最有效的",
        "P4：AI 生成有随机性，'像巫术一样'（E2 原话）",
        "P5：所有人都想要控制权——自动完成不等于好体验",
    ]
    add_bullet(slide, pain_pts, Inches(0.55), Inches(4.4), Inches(5.7), Inches(2.4),
               font_size=10, color=C_DARK)

    # 右：User Study
    add_rect(slide, Inches(6.7), Inches(1.0), Inches(6.2), Inches(6.1), C_WHITE, C_PURPLE, Pt(1.5))
    add_rect(slide, Inches(6.7), Inches(1.0), Inches(6.2), Inches(0.5), C_PURPLE)
    add_text(slide, "User Study（15名终端用户）", Inches(6.85), Inches(1.05), Inches(5.9), Inches(0.4),
             font_size=13, bold=True, color=C_WHITE)

    add_text(slide, "研究设计", Inches(6.85), Inches(1.62), Inches(5.9), Inches(0.3),
             font_size=11, bold=True, color=C_DARK2)
    study = [
        "参与者：15 人，使用自己的私人照片进行测试",
        "任务：用 Imago Obscura 分析并决定是否/如何发布图片",
        "方法：实验室研究 + 访谈 + 问卷",
        "关键指标：隐私风险意识提升、风险规避能力、明智决策能力",
    ]
    add_bullet(slide, study, Inches(6.85), Inches(1.92), Inches(5.9), Inches(1.55),
               font_size=10, color=C_DARK)

    add_text(slide, "✅ 积极发现", Inches(6.85), Inches(3.55), Inches(5.9), Inches(0.3),
             font_size=11, bold=True, color=C_GREEN)
    pos = [
        "用户喜欢表达自己担忧的过程——这帮助他们聚焦最在意的风险",
        "全面呈现风险 + 推荐脱敏方案 → 支持更明智的决策",
        "一键脱敏大大降低技术门槛",
        "保留精细控制权 → 用户感到安全和高自主性",
    ]
    add_bullet(slide, pos, Inches(6.85), Inches(3.85), Inches(5.9), Inches(1.55),
               font_size=10, color=C_DARK)

    add_text(slide, "⚠️ 局限与风险", Inches(6.85), Inches(5.48), Inches(5.9), Inches(0.3),
             font_size=11, bold=True, color=C_RED)
    limits = [
        "恶意使用风险：降低了制作虚假信息的门槛（deepfake）",
        "缺乏 Guardrails（防护栏）防止滥用",
        "AI 输出的随机性仍可能导致用户困惑",
    ]
    add_bullet(slide, limits, Inches(6.85), Inches(5.78), Inches(5.9), Inches(1.3),
               font_size=10, color=C_DARK)

    return slide

# ================================================================
def slide_zhao_detail(prs):
    slide = bg(prs, C_BG)
    """论文2 — Jun Zhao 详解"""
    section_header(slide, "论文 2  |  Parents & Children's Online Privacy — 用户认知研究", "② 用户理解")

    # 左：基本信息
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(5.9), Inches(6.1), C_WHITE, C_AMBER, Pt(1.5))
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(5.9), Inches(0.5), C_AMBER)
    add_text(slide, "基本信息", Inches(0.55), Inches(1.05), Inches(5.6), Inches(0.4),
             font_size=14, bold=True, color=C_WHITE)

    info = [
        ("论文", "Are Children Well-Supported by Their Parents Concerning Online Privacy Risks, and Who Supports the Parents?"),
        ("作者", "Jun Zhao"),
        ("发表", "Proc. of Privacy & Trust, 2018"),
        ("研究对象", "6-11 岁儿童的父母"),
        ("方法", "访谈 + 问卷调查"),
        ("核心问题", "父母是否意识到孩子在线隐私风险？他们如何应对？谁在支持父母？"),
    ]
    y = Inches(1.62)
    for label, value in info:
        add_text(slide, f"{label}：", Inches(0.55), y, Inches(1.0), Inches(0.35),
                 font_size=11, bold=True, color=C_AMBER)
        add_text(slide, value, Inches(1.55), y, Inches(4.6), Inches(0.7),
                 font_size=11, color=C_DARK)
        y += Inches(0.72)

    # 右：关键发现
    add_rect(slide, Inches(6.6), Inches(1.0), Inches(6.4), Inches(6.1), C_WHITE, C_AMBER, Pt(1.5))
    add_rect(slide, Inches(6.6), Inches(1.0), Inches(6.4), Inches(0.5), C_AMBER)
    add_text(slide, "关键研究发现（4个核心发现）", Inches(6.75), Inches(1.05), Inches(6.1), Inches(0.4),
             font_size=13, bold=True, color=C_WHITE)

    findings = [
        ("F1", "认知偏差",
         "父母往往认为孩子太小、不需要担心在线隐私问题——但实际上孩子已经面临真实风险。\n"
         "他们低估了数字世界中信息传播的速度和范围。"),
        ("F2", "保护性策略 > 教育性策略",
         "父母主要采取「保护性」策略：限制使用、监控活动、禁止某些操作。\n"
         "而非「教育性」策略：教会孩子识别风险、自主判断。\n"
         "结果：孩子仍然缺乏隐私保护意识和能力。"),
        ("F3", "信息不对称",
         "父母自认为了解风险，但对具体威胁的认知不足。\n"
         "不清楚 sharenting 可能带来的长期后果（未来雇主/大学/恋人的审视）。"),
        ("F4", "工具缺乏",
         "缺乏有效的工具帮助父母在「分享」和「保护」之间找到平衡。\n"
         "现有工具要么太技术化，要么只针对单一风险点。"),
    ]
    y = Inches(1.62)
    for fid, title, desc in findings:
        add_rect(slide, Inches(6.7), y, Inches(6.2), Inches(1.45), RGBColor(0xFF, 0xF8, 0xF0))
        add_rect(slide, Inches(6.7), y, Inches(0.55), Inches(1.45), C_AMBER)
        add_text(slide, fid, Inches(6.7), y + Inches(0.5), Inches(0.55), Inches(0.4),
                 font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, Inches(7.35), y + Inches(0.08), Inches(5.4), Inches(0.35),
                 font_size=12, bold=True, color=C_AMBER)
        add_text(slide, desc, Inches(7.35), y + Inches(0.45), Inches(5.4), Inches(0.95),
                 font_size=10, color=C_DARK)
        y += Inches(1.55)

    return slide

# ================================================================
def slide_vlm_detail(prs):
    slide = bg(prs, C_BG)
    """论文3 — VLM Content Moderation 详解"""
    section_header(slide, "论文 3  |  VLM Content Moderation — 多模态内容审核技术", "③ 技术路线")

    # 左：基本信息
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(5.9), Inches(6.1), C_WHITE, C_PINK, Pt(1.5))
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(5.9), Inches(0.5), C_PINK)
    add_text(slide, "基本信息", Inches(0.55), Inches(1.05), Inches(5.6), Inches(0.4),
             font_size=14, bold=True, color=C_WHITE)

    info = [
        ("论文", "The Potential of Vision-Language Models for Content Moderation of Children's Videos"),
        ("作者", "Syed Hammad Ahmed, Shengnan Hu, Gita Sukthankar"),
        ("发表", "arXiv 2023"),
        ("补充", "Enhanced Multimodal Content Moderation of Children's Videos (arXiv 2024) — 音频+视觉融合"),
    ]
    y = Inches(1.62)
    for label, value in info:
        add_text(slide, f"{label}：", Inches(0.55), y, Inches(1.0), Inches(0.35),
                 font_size=11, bold=True, color=C_PINK)
        add_text(slide, value, Inches(1.55), y, Inches(4.6), Inches(0.7),
                 font_size=11, color=C_DARK)
        y += Inches(0.72)

    add_text(slide, "核心挑战", Inches(0.55), Inches(4.15), Inches(5.6), Inches(0.3),
             font_size=12, bold=True, color=C_DARK2)
    challenges = [
        "内容审核不只判断「有没有暴力/色情」",
        "很多不当内容以「教育内容」形式伪装（如垃圾广告伪装成儿童动画）",
        " scammers 利用视觉相似性绕过审核",
        "2024 年升级版：纯视觉无法检测「听觉上的不当内容」（恐怖音效/不适宜音乐）",
    ]
    add_bullet(slide, challenges, Inches(0.55), Inches(4.45), Inches(5.6), Inches(2.2),
               font_size=10, color=C_DARK)

    # 右：技术方案
    add_rect(slide, Inches(6.6), Inches(1.0), Inches(6.4), Inches(6.1), C_WHITE, C_PINK, Pt(1.5))
    add_rect(slide, Inches(6.6), Inches(1.0), Inches(6.4), Inches(0.5), C_PINK)
    add_text(slide, "技术方案：VLM 零样本内容审核", Inches(6.75), Inches(1.05), Inches(6.1), Inches(0.4),
             font_size=13, bold=True, color=C_WHITE)

    tech_pts = [
        ("方法", "使用 Vision-Language Models（VLM）进行零样本检测，不需要针对每个类别训练"),
        ("核心思路", "用自然语言文字提示（text prompts）引导模型识别不当内容"),
        ("优势", "不需要大量「儿童不宜内容」的标注训练数据，减少数据收集的伦理问题"),
        ("可检测类型", "暴力、色情等明显内容之外，还能检测细微的不当内容（如伪装成教育内容的垃圾视频）"),
        ("2024 升级", "Audiovisual Fusion — 音频+视觉多模态融合：视频画面正常但音频含恐怖音效也能被检测"),
    ]
    y = Inches(1.62)
    for label, desc in tech_pts:
        add_rect(slide, Inches(6.7), y, Inches(6.2), Inches(0.9), C_LIGHT if tech_pts.index((label,desc))%2==0 else C_WHITE)
        add_text(slide, label, Inches(6.85), y + Inches(0.08), Inches(1.0), Inches(0.3),
                 font_size=10, bold=True, color=C_PINK)
        add_text(slide, desc, Inches(7.85), y + Inches(0.08), Inches(4.9), Inches(0.75),
                 font_size=10, color=C_DARK)
        y += Inches(0.96)

    add_text(slide, "对 Future Child Posting 的直接启示",
             Inches(6.75), Inches(6.0), Inches(6.1), Inches(0.3),
             font_size=11, bold=True, color=C_DARK2)
    add_text(slide,
             "VLM 可以识别图片中的场景（公园/医院/学校/家庭内部）而不仅是独立对象。"
             "文字提示（prompts）可以引导模型关注特定的隐私风险类别。"
             "结合 Imago Obscura 的框架，这是图片分析模块的核心技术选型。",
             Inches(6.75), Inches(6.3), Inches(6.1), Inches(0.75),
             font_size=10, color=C_DARK)

    return slide

# ================================================================
def slide_risks_for_children(prs):
    slide = bg(prs, C_BG)
    """Sharenting 特有的儿童隐私风险（更具体的）"""
    section_header(slide, "延伸分析  |  儿童内容的特殊隐私风险", "Special Risks")

    add_text(slide, "Sharenting 的核心矛盾：父母有分享权，但孩子没有同意权，且后果由孩子承担",
             Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.45),
             font_size=14, bold=True, color=C_DARK2)

    # 三列
    cols = [
        {
            "title": "信息层面风险",
            "color": C_PURPLE,
            "items": [
                "🏫 学校信息：校服标志、书包、班级活动 → 可定位孩子学校",
                "🏠 家庭位置：窗外地标、门牌、快递单 → 可定位住址",
                "👶 身体信息：裸浴/如厕训练照片 → 许多国家的法律红线",
                "🏥 医疗信息：医院照片、疫苗记录 → 健康数据泄露",
                "🎂 生日信息：具体日期+姓名 → 身份盗窃原材料",
                "💊 健康状况：特殊饮食/药物照片 → 健康记录推断",
            ],
        },
        {
            "title": "数字足迹风险",
            "color": C_AMBER,
            "items": [
                "⏳ 终身影响：一旦发布终身无法撤回，孩子无法拒绝",
                "👔 未来审视：未来雇主/大学/恋人可能搜索到这些内容",
                "🤖 AI 抓取：照片可能被用于 AI 训练数据",
                "📊 数据画像：在不知情的情况下形成孩子的数字画像",
                "🔗 跨平台传播：截图扩散到其他平台，完全脱离控制",
                "🕵️ 关系推断：点赞/评论关系 → 社交网络暴露",
            ],
        },
        {
            "title": "心理与伦理风险",
            "color": C_PINK,
            "items": [
                "😔 尴尬未来：孩子长大后对自己「数字身份」失去控制",
                "🆘 情感伤害：被公开讨论、嘲笑、欺凌的风险",
                "⚖️ Consent 问题：孩子没有机会同意或拒绝",
                "🤝 旁观者权利：照片中其他孩子家长的同意权",
                "💔 离婚风险：共享账号内容在离婚时的法律归属",
                "📵 上瘾机制：点赞驱动过度分享的心理动因",
            ],
        },
    ]

    x = Inches(0.4)
    for col in cols:
        w = Inches(4.18)
        add_rect(slide, x, Inches(1.65), w, Inches(5.4), C_WHITE, col["color"], Pt(2))
        add_rect(slide, x, Inches(1.65), w, Inches(0.55), col["color"])
        add_text(slide, col["title"], x + Inches(0.15), Inches(1.72), w - Inches(0.3), Inches(0.42),
                 font_size=13, bold=True, color=C_WHITE)
        y = Inches(2.35)
        for item in col["items"]:
            add_text(slide, item, x + Inches(0.15), y, w - Inches(0.3), Inches(0.5),
                     font_size=11, color=C_DARK)
            y += Inches(0.6)
        x += Inches(4.43)

    return slide

# ================================================================
def slide_project_rq(prs):
    slide = bg(prs, C_BG)
    """我们项目的 RQ + 技术方案"""
    section_header(slide, "项目定位  |  RQ 与技术方案", "Project")

    # RQ Box
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(12.5), Inches(1.1), C_LIGHT, C_PURPLE, Pt(2))
    add_text(slide, "RQ：AI 辅助的『发布前风险扫描』能否有效帮助父母识别和规避分享孩子内容时的隐私和伦理风险？",
             Inches(0.6), Inches(1.12), Inches(12.2), Inches(0.8),
             font_size=15, bold=True, color=C_DARK2)

    # 左：功能设计
    add_rect(slide, Inches(0.4), Inches(2.3), Inches(6.0), Inches(4.8), C_WHITE, C_PURPLE, Pt(1.5))
    add_rect(slide, Inches(0.4), Inches(2.3), Inches(6.0), Inches(0.5), C_PURPLE)
    add_text(slide, "MVP 功能设计（参考 Imago Obscura DR1-5）", Inches(0.55), Inches(2.35), Inches(5.7), Inches(0.4),
             font_size=13, bold=True, color=C_WHITE)

    features = [
        ("📝 文字输入", "用户输入想发的文案或描述分享意图"),
        ("🖼️ 图片上传", "上传待检查的图片（可多张）"),
        ("🔍 AI 风险扫描", "识别5类隐私风险 + 给出风险等级"),
        ("💡 脱敏建议", "对每类风险给出文字建议（参考 Imago Obscura）"),
        ("✍️ 改写文案", "AI 生成脱敏后的替代文案"),
        ("📊 风险报告", "汇总所有检测到的风险点和总体评分"),
    ]
    y = Inches(2.95)
    for feat, desc in features:
        add_rect(slide, Inches(0.55), y, Inches(5.7), Inches(0.62), C_LIGHT)
        add_text(slide, feat, Inches(0.7), y + Inches(0.08), Inches(1.5), Inches(0.3),
                 font_size=11, bold=True, color=C_PURPLE)
        add_text(slide, desc, Inches(2.2), y + Inches(0.1), Inches(3.9), Inches(0.45),
                 font_size=11, color=C_DARK)
        y += Inches(0.68)

    # 右：技术方案
    add_rect(slide, Inches(6.7), Inches(2.3), Inches(6.2), Inches(4.8), C_WHITE, C_TEAL, Pt(1.5))
    add_rect(slide, Inches(6.7), Inches(2.3), Inches(6.2), Inches(0.5), C_TEAL)
    add_text(slide, "技术方案（初稿）", Inches(6.85), Inches(2.35), Inches(5.9), Inches(0.4),
             font_size=13, bold=True, color=C_WHITE)

    tech = [
        ("文字扫描", "LLM API（GPT-4o / Claude 等）+ NER 实体识别 + 规则匹配"),
        ("图片扫描", "VLM（GPT-4V 或等效多模态模型）+ 场景分类"),
        ("风险检测", "基于 Imago Obscura 5类风险构建的规则 + LLM 综合判断"),
        ("前端", "网页 App（Next.js 或纯前端）"),
        ("后端", "Spring Boot 或 FastAPI（轻量级）"),
        ("API 设计", "POST /analyze（文字+图片）→ 返回风险报告 JSON"),
    ]
    y = Inches(2.95)
    for label, desc in tech:
        add_rect(slide, Inches(6.8), y, Inches(6.0), Inches(0.62), C_LIGHT)
        add_text(slide, label, Inches(6.95), y + Inches(0.08), Inches(1.3), Inches(0.3),
                 font_size=11, bold=True, color=C_TEAL)
        add_text(slide, desc, Inches(8.3), y + Inches(0.08), Inches(4.4), Inches(0.5),
                 font_size=11, color=C_DARK)
        y += Inches(0.68)

    return slide

# ================================================================
def slide_mvp_overview(prs):
    """MVP 原型概览"""
    slide = bg(prs, C_BG)
    section_header(slide, "MVP 原型  |  系统概览", "Prototype")

    # 标题说明
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.7), C_LIGHT, C_PURPLE, Pt(1.5))
    add_text(slide, "基于 Imago Obscura 风险框架，构建了一个可演示的 MVP 原型，支持图片+文字联合分析，识别 5 类儿童隐私风险并给出脱敏建议。",
             Inches(0.6), Inches(1.05), Inches(12.2), Inches(0.6),
             font_size=13, color=C_DARK2)

    # 三个核心指标卡片
    cards = [
        ("🖼️ 图片支持", "最多 9 张图\n联合分析", C_PURPLE),
        ("🔍 风险类别", "R1–R5 共 5 类\n隐私风险检测", C_AMBER),
        ("💡 脱敏建议", "逐条风险\n修改建议", C_GREEN),
    ]
    x = Inches(0.4)
    for title, desc, color in cards:
        add_rect(slide, x, Inches(1.9), Inches(4.0), Inches(2.0), C_WHITE, color, Pt(2))
        add_rect(slide, x, Inches(1.9), Inches(4.0), Inches(0.55), color)
        add_text(slide, title, x + Inches(0.15), Inches(1.97), Inches(3.7), Inches(0.42),
                 font_size=14, bold=True, color=C_WHITE)
        add_text(slide, desc, x + Inches(0.15), Inches(2.55), Inches(3.7), Inches(1.2),
                 font_size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
        x += Inches(4.3)

    # 工作流程
    add_text(slide, "用户使用流程", Inches(0.4), Inches(4.15), Inches(3.0), Inches(0.4),
             font_size=13, bold=True, color=C_DARK)
    flow = ["上传图片 + 输入文案", "AI 风险扫描", "生成报告 + 脱敏建议"]
    colors_flow = [C_PURPLE, C_AMBER, C_GREEN]
    x = Inches(0.4)
    for i, step in enumerate(flow):
        add_rect(slide, x, Inches(4.6), Inches(3.5), Inches(0.8), colors_flow[i])
        add_text(slide, step, x + Inches(0.1), Inches(4.72), Inches(3.3), Inches(0.55),
                 font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            add_text(slide, "→", x + Inches(3.55), Inches(4.72), Inches(0.5), Inches(0.55),
                     font_size=20, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        x += Inches(4.0)

    # 技术栈
    add_text(slide, "技术栈", Inches(0.4), Inches(5.65), Inches(2.0), Inches(0.4),
             font_size=13, bold=True, color=C_DARK)
    tech_stack = [
        ("后端", "FastAPI + Python", C_TEAL),
        ("前端", "Next.js + React", C_PURPLE),
        ("AI 模型", "qwen3.6-plus（阿里云百炼）", C_AMBER),
        ("代理", "Clash Verge (7897)", C_GRAY),
    ]
    x = Inches(0.4)
    for label, val, color in tech_stack:
        add_rect(slide, x, Inches(6.1), Inches(3.1), Inches(0.95), color)
        add_text(slide, label, x + Inches(0.1), Inches(6.15), Inches(2.9), Inches(0.35),
                 font_size=10, bold=True, color=C_WHITE)
        add_text(slide, val, x + Inches(0.1), Inches(6.48), Inches(2.9), Inches(0.45),
                 font_size=11, bold=True, color=C_WHITE)
        x += Inches(3.28)

    return slide


def slide_mvp_upload(prs):
    """MVP 上传界面"""
    slide = bg(prs, C_BG)
    section_header(slide, "MVP 原型  |  上传与分析界面", "Demo")

    # 左侧：界面描述
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(6.0), Inches(6.2), C_WHITE, C_PURPLE, Pt(1.5))
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(6.0), Inches(0.55), C_PURPLE)
    add_text(slide, "上传与输入面板", Inches(0.55), Inches(1.0), Inches(5.7), Inches(0.42),
             font_size=13, bold=True, color=C_WHITE)

    features = [
        ("📷 图片上传区", "支持拖拽或点击上传最多 9 张图片，显示缩略图预览，提示文件大小和格式限制"),
        ("📝 文字输入框", "用户输入准备发布的内容描述或分享意图，支持多行文本输入"),
        ("🚀 一键分析", "点击「开始分析」按钮，触发 AI 风险扫描，显示加载动画"),
        ("⏳ 等待过程", "实时显示分析进度和状态，给出预估时间提示"),
    ]
    y = Inches(1.65)
    for title, desc in features:
        add_rect(slide, Inches(0.55), y, Inches(5.7), Inches(1.2), C_LIGHT)
        add_text(slide, title, Inches(0.7), y + Inches(0.1), Inches(5.4), Inches(0.38),
                 font_size=12, bold=True, color=C_PURPLE)
        add_text(slide, desc, Inches(0.7), y + Inches(0.48), Inches(5.4), Inches(0.65),
                 font_size=11, color=C_DARK)
        y += Inches(1.28)

    # 右侧：交互说明
    add_rect(slide, Inches(6.7), Inches(0.95), Inches(6.2), Inches(6.2), C_WHITE, C_AMBER, Pt(1.5))
    add_rect(slide, Inches(6.7), Inches(0.95), Inches(6.2), Inches(0.55), C_AMBER)
    add_text(slide, "交互设计细节", Inches(6.85), Inches(1.0), Inches(5.9), Inches(0.42),
             font_size=13, bold=True, color=C_WHITE)

    details = [
        ("无压力设计", "用户可以随时修改图片或文字，多次提交分析而不产生实际发布行为"),
        ("隐私保护", "图片仅在分析过程中使用，不会持久化存储，消除用户隐私顾虑"),
        ("渐进式引导", "为首次使用的家长提供操作提示，降低使用门槛"),
        ("多图协同", "多张图片之间可标注关联（如「左图是孩子，右图是家庭成员」），帮助 AI 理解场景"),
    ]
    y = Inches(1.65)
    for title, desc in details:
        add_rect(slide, Inches(6.8), y, Inches(6.0), Inches(1.2), C_LIGHT)
        add_text(slide, title, Inches(6.95), y + Inches(0.1), Inches(5.7), Inches(0.38),
                 font_size=12, bold=True, color=C_AMBER)
        add_text(slide, desc, Inches(6.95), y + Inches(0.48), Inches(5.7), Inches(0.65),
                 font_size=11, color=C_DARK)
        y += Inches(1.28)

    return slide


def slide_mvp_report(prs):
    """MVP 风险报告页面"""
    slide = bg(prs, C_BG)
    section_header(slide, "MVP 原型  |  风险报告展示", "Demo")

    # 顶部说明
    add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.65), C_LIGHT, C_PINK, Pt(1.5))
    add_text(slide, "分析完成后，展示 R1–R5 每类风险的检测结果、严重程度（低/中/高）以及针对性的脱敏修改建议。",
             Inches(0.6), Inches(1.05), Inches(12.2), Inches(0.5),
             font_size=13, color=C_DARK2)

    # 5类风险卡片
    risks = [
        ("R1", "自我披露风险", "检测孩子外貌、姓名、学校等可识别信息披露", C_RED),
        ("R2", "背景泄露风险", "家庭住址、旅行行程、活动场所等背景信息", C_AMBER),
        ("R3", "行为风险", "危险行为、不当内容或不安全情境的展示", C_PINK),
        ("R4", "元数据风险", "照片 EXIF 地理位置、时间戳等隐藏信息的泄露", C_PURPLE),
        ("R5", "二次传播风险", "内容被截图、转发或用于非预期目的的可能性", C_TEAL),
    ]
    x = Inches(0.4)
    for code, title, desc, color in risks:
        add_rect(slide, x, Inches(1.8), Inches(2.45), Inches(2.8), C_WHITE, color, Pt(2))
        add_rect(slide, x, Inches(1.8), Inches(2.45), Inches(0.5), color)
        add_text(slide, code, x + Inches(0.1), Inches(1.85), Inches(1.0), Inches(0.42),
                 font_size=16, bold=True, color=C_WHITE)
        add_text(slide, title, x + Inches(1.1), Inches(1.88), Inches(1.25), Inches(0.38),
                 font_size=9, bold=True, color=C_WHITE)
        add_text(slide, desc, x + Inches(0.1), Inches(2.38), Inches(2.25), Inches(2.0),
                 font_size=10, color=C_DARK)
        x += Inches(2.55)

    # 报告UI说明
    add_rect(slide, Inches(0.4), Inches(4.85), Inches(7.5), Inches(2.3), C_WHITE, C_GREEN, Pt(1.5))
    add_rect(slide, Inches(0.4), Inches(4.85), Inches(7.5), Inches(0.5), C_GREEN)
    add_text(slide, "风险报告 UI 设计", Inches(0.55), Inches(4.9), Inches(7.2), Inches(0.42),
             font_size=12, bold=True, color=C_WHITE)
    report_features = [
        "• 每类风险独立卡片，带颜色编码（红/橙/粉/紫/青）",
        "• 风险严重程度：低（L）/ 中（M）/ 高（H）三级标注",
        "• 每条风险配有「为什么会这样」的解释文字",
        "• 「修改建议」按钮可展开 AI 给出的脱敏方案",
    ]
    y = Inches(5.45)
    for item in report_features:
        add_text(slide, item, Inches(0.6), y, Inches(7.1), Inches(0.38),
                 font_size=11, color=C_DARK)
        y += Inches(0.4)

    # AI改写说明
    add_rect(slide, Inches(8.2), Inches(4.85), Inches(4.7), Inches(2.3), C_WHITE, C_TEAL, Pt(1.5))
    add_rect(slide, Inches(8.2), Inches(4.85), Inches(4.7), Inches(0.5), C_TEAL)
    add_text(slide, "AI 文案改写功能", Inches(8.35), Inches(4.9), Inches(4.4), Inches(0.42),
             font_size=12, bold=True, color=C_WHITE)
    rewrite = [
        "• 「生成安全文案」：AI 输出一整段脱敏后的替代文案",
        "• 保留原意，过滤可识别信息",
        "• 用户可直接复制使用，或继续调整",
    ]
    y = Inches(5.45)
    for item in rewrite:
        add_text(slide, item, Inches(8.35), y, Inches(4.4), Inches(0.45),
                 font_size=11, color=C_DARK)
        y += Inches(0.45)

    return slide


def slide_mvp_tech(prs):
    """MVP 技术实现"""
    slide = bg(prs, C_BG)
    section_header(slide, "MVP 原型  |  技术实现与数据流", "Tech")

    # 数据流图（文字版）
    add_text(slide, "系统数据流", Inches(0.4), Inches(0.95), Inches(3.0), Inches(0.4),
             font_size=13, bold=True, color=C_DARK)

    flow_steps = [
        ("用户", "上传图片 + 文字", C_PURPLE),
        ("前端", "Base64 编码", C_AMBER),
        ("后端", "FastAPI 接收", C_TEAL),
        ("AI 服务", "qwen3.6-plus", C_PINK),
        ("返回", "风险报告 JSON", C_GREEN),
        ("前端", "渲染报告 UI", C_PURPLE),
    ]
    x = Inches(0.4)
    for i, (node, action, color) in enumerate(flow_steps):
        add_rect(slide, x, Inches(1.4), Inches(2.0), Inches(0.85), color)
        add_text(slide, node, x + Inches(0.1), Inches(1.45), Inches(1.8), Inches(0.38),
                 font_size=12, bold=True, color=C_WHITE)
        add_text(slide, action, x + Inches(0.1), Inches(1.8), Inches(1.8), Inches(0.38),
                 font_size=9, color=C_WHITE)
        if i < len(flow_steps) - 1:
            add_text(slide, "→", x + Inches(2.05), Inches(1.55), Inches(0.3), Inches(0.5),
                     font_size=18, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        x += Inches(2.35)

    # API 设计
    add_rect(slide, Inches(0.4), Inches(2.5), Inches(6.0), Inches(4.6), C_WHITE, C_PURPLE, Pt(1.5))
    add_rect(slide, Inches(0.4), Inches(2.5), Inches(6.0), Inches(0.55), C_PURPLE)
    add_text(slide, "API 设计（FastAPI）", Inches(0.55), Inches(2.56), Inches(5.7), Inches(0.42),
             font_size=13, bold=True, color=C_WHITE)

    api_items = [
        ("POST /api/analyze", "接收 { images: [base64], text: string }，返回风险报告"),
        ("请求示例", '{ "images": ["base64..."], "text": "宝宝今天在公园玩得很开心！" }'),
        ("响应结构", '{ risks: [{ category, level, description, suggestion }] }'),
        ("错误处理", "图片超限 / API 超时 / 格式错误均有友好提示"),
    ]
    y = Inches(3.15)
    for title, desc in api_items:
        add_rect(slide, Inches(0.55), y, Inches(5.7), Inches(0.95), C_LIGHT)
        add_text(slide, title, Inches(0.7), y + Inches(0.1), Inches(5.4), Inches(0.35),
                 font_size=11, bold=True, color=C_PURPLE)
        add_text(slide, desc, Inches(0.7), y + Inches(0.45), Inches(5.4), Inches(0.45),
                 font_size=10, color=C_DARK)
        y += Inches(1.02)

    # 风险检测 prompt
    add_rect(slide, Inches(6.7), Inches(2.5), Inches(6.2), Inches(4.6), C_WHITE, C_AMBER, Pt(1.5))
    add_rect(slide, Inches(6.7), Inches(2.5), Inches(6.2), Inches(0.55), C_AMBER)
    add_text(slide, "AI 风险检测 Prompt（核心逻辑）", Inches(6.85), Inches(2.56), Inches(5.9), Inches(0.42),
             font_size=13, bold=True, color=C_WHITE)

    prompt_items = [
        ("System Prompt", "你是一个儿童隐私安全专家，参考 Imago Obscura 的 5 类风险框架，对用户内容进行隐私风险分析..."),
        ("R1 检测", "孩子外貌、姓名、年龄、学校等可识别信息"),
        ("R2 检测", "家庭住址、地理位置、旅游行程等背景信息"),
        ("R3 检测", "危险行为、不安全场景、过度暴露身体等"),
        ("R4 检测", "EXIF 地理位置、时间戳、GPS 等元数据"),
        ("R5 检测", "内容被二次传播或滥用的可能性评估"),
    ]
    y = Inches(3.15)
    for label, desc in prompt_items:
        add_rect(slide, Inches(6.8), y, Inches(6.0), Inches(0.68), C_LIGHT)
        add_text(slide, label, Inches(6.95), y + Inches(0.08), Inches(1.5), Inches(0.3),
                 font_size=10, bold=True, color=C_AMBER)
        add_text(slide, desc, Inches(8.45), y + Inches(0.08), Inches(4.2), Inches(0.55),
                 font_size=10, color=C_DARK)
        y += Inches(0.74)

    return slide


# ================================================================
def slide_next_steps(prs):
    slide = bg(prs, C_BG)
    """下一步计划"""
    section_header(slide, "下一步  |  Next Steps", "Roadmap")

    steps = [
        {
            "phase": "Week 1-2", "title": "设计研究",
            "color": C_PURPLE,
            "items": [
                "深度访谈 5-8 位家长，了解 sharenting 现状和痛点",
                "确定目标用户画像（孩子年龄段、分享频率、平台偏好）",
                "明确 MVP 功能边界（先做文字+单图扫描）",
                "起草系统工作流和交互流程",
            ],
        },
        {
            "phase": "Week 3-4", "title": "原型开发",
            "color": C_AMBER,
            "items": [
                "接入 LLM API，实现文字隐私风险扫描",
                "接入 VLM，实现图片场景和敏感内容识别",
                "开发风险报告 UI（网页端）",
                "迭代优化 AI 提示词，提升检测准确率",
            ],
        },
        {
            "phase": "Week 5-7", "title": "用户研究",
            "color": C_PINK,
            "items": [
                "招募 10-15 名真实家长",
                "任务测试：使用原型完成发布前风险扫描",
                "访谈 + 问卷（SUS + 信任度 + 接受度）",
                "数据分析 → 验证 RQ",
            ],
        },
        {
            "phase": "Week 8+", "title": "论文产出",
            "color": C_TEAL,
            "items": [
                "整理研究发现，撰写论文",
                "目标投稿：CHI / CSCW / IDC",
                "（可选）开源原型代码和数据集",
            ],
        },
    ]

    x = Inches(0.4)
    for step in steps:
        w = Inches(3.1)
        add_rect(slide, x, Inches(1.0), w, Inches(6.1), C_WHITE, step["color"], Pt(2))
        add_rect(slide, x, Inches(1.0), w, Inches(0.55), step["color"])
        add_text(slide, step["phase"], x + Inches(0.12), Inches(1.07), Inches(1.1), Inches(0.42),
                 font_size=12, bold=True, color=C_WHITE)
        add_text(slide, step["title"], x + Inches(1.25), Inches(1.12), w - Inches(1.4), Inches(0.38),
                 font_size=12, bold=True, color=C_WHITE)
        y = Inches(1.7)
        for item in step["items"]:
            add_text(slide, "→ " + item, x + Inches(0.12), y, w - Inches(0.24), Inches(0.65),
                     font_size=11, color=C_DARK)
            y += Inches(0.72)
        x += Inches(3.28)

    return slide

# ================================================================
def main():
    prs = new_prs()
    slide_cover(prs)
    slide_problem(prs)
    slide_papers_overview(prs)
    slide_imago_threat(prs)
    slide_imago_dr(prs)
    slide_imago_tech(prs)
    slide_imago_user_study(prs)
    slide_zhao_detail(prs)
    slide_vlm_detail(prs)
    slide_risks_for_children(prs)
    slide_project_rq(prs)
    slide_mvp_overview(prs)
    slide_mvp_upload(prs)
    slide_mvp_report(prs)
    slide_mvp_tech(prs)
    slide_next_steps(prs)

    out = "/Users/sunxiaoxuan/Desktop/FutureChildPosting_Research.pptx"
    prs.save(out)
    print(f"✅ 详细版 PPT 已保存：{out}")

if __name__ == "__main__":
    main()